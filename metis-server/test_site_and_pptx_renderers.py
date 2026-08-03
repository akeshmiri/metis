"""
§12.5 Site renderer + §4.6.1 PPTX renderer -- both real, both call the
same metis_mcp/academy.py content-assembly stage (REQ-METIS-ACD-07), both
tested against real output files, not mocked rendering.
"""
import os
import shutil
import sys
import tempfile
from pathlib import Path

from neo4j import GraphDatabase
from pptx import Presentation

from metis_mcp.site_renderer import render_site, render_test_design_report
from metis_mcp.pptx_renderer import render_quality_deck
from demo_data.generate_demo_data import generate, wipe_demo_data, Scale

NEO4J_URI = os.environ.get("METIS_NEO4J_URI", "bolt://localhost:7687")
NEO4J_USER = os.environ.get("METIS_NEO4J_USER", "neo4j")
NEO4J_PASSWORD = os.environ.get("METIS_NEO4J_PASSWORD")

_driver = None


def _session():
    global _driver
    if _driver is None:
        _driver = GraphDatabase.driver(NEO4J_URI, auth=(NEO4J_USER, NEO4J_PASSWORD))
    return _driver.session()


def test_render_site_writes_real_html_for_every_academy_page():
    with tempfile.TemporaryDirectory() as d:
        with _session() as s:
            result = render_site(d, session=s)
        assert os.path.isfile(result["index"])
        assert len(result["academy_pages"]) == 4
        for path in result["academy_pages"]:
            assert os.path.isfile(path)
            with open(path, encoding="utf-8") as f:
                content = f.read()
            assert "<html>" in content
            assert "REPLACE" not in content, "no leftover placeholder text"
            assert "Academy content version: 1.0.0" in content, \
                "REQ-METIS-ACD-06: version must actually be visible, not just computed"

        with open(result["index"], encoding="utf-8") as f:
            index_content = f.read()
        assert "quality_score" in index_content.lower() or "not computable" in index_content

        # REQ-METIS-ACD-05: a real changelog page, sourced from real
        # :Constitution revision history (populated by this session's own
        # earlier `python3 -m metis_mcp.constitution_gate` run).
        assert os.path.isfile(result["changelog"])
        with open(result["changelog"], encoding="utf-8") as f:
            changelog_content = f.read()
        assert "CONST-047" in changelog_content
        assert "recorded (initial state)" in changelog_content or "changed from" in changelog_content


def test_render_site_index_links_resolve_to_real_files():
    with tempfile.TemporaryDirectory() as d:
        result = render_site(d)  # no session -- index still renders, just skips the live snapshot
        with open(result["index"], encoding="utf-8") as f:
            index_content = f.read()
        for path in result["academy_pages"]:
            rel = "academy/" + os.path.basename(path)
            assert rel in index_content


def test_render_test_design_report_shows_real_login_example_techniques():
    """Session 10: 'anyone can go and check this model' -- the Site page,
    not just the MCP tool, must show real technique names from the real
    Intent/TestDesign backbone."""
    wipe_demo_data(NEO4J_URI, NEO4J_USER, NEO4J_PASSWORD)
    generate(NEO4J_URI, NEO4J_USER, NEO4J_PASSWORD, scale=Scale(factor=0.05), seed=7)
    try:
        with tempfile.TemporaryDirectory() as d:
            with _session() as s:
                path = render_test_design_report(Path(d), session=s,
                                                  scope={"requirement_id": "demo:login:requirement:t3-lockout"})
            assert os.path.isfile(path)
            with open(path, encoding="utf-8") as f:
                content = f.read()
            assert "Boundary Value Analysis" in content
            assert "unit" in content and "api_functional" in content
    finally:
        wipe_demo_data(NEO4J_URI, NEO4J_USER, NEO4J_PASSWORD)


def test_render_quality_deck_produces_a_real_reopenable_pptx():
    with tempfile.TemporaryDirectory() as d:
        output_path = os.path.join(d, "quality-snapshot.pptx")
        with _session() as s:
            result = render_quality_deck(s, output_path)
        assert os.path.isfile(output_path)
        assert result["slide_count"] >= 2
        assert result["file_qa_passed"] is True
        assert result["content_qa_passed"] is True

        # Independently re-open with python-pptx directly, not just trusting the module's own report.
        prs = Presentation(output_path)
        assert len(prs.slides) == result["slide_count"]
        all_text = "\n".join(
            shape.text_frame.text for slide in prs.slides for shape in slide.shapes if shape.has_text_frame
        )
        assert "quality_score" in all_text.lower()
        assert result["generated_at"] in all_text, "the point-in-time anchor must actually be on a slide"


if __name__ == "__main__":
    if not NEO4J_PASSWORD:
        print("METIS_NEO4J_PASSWORD is not set.", file=sys.stderr)
        sys.exit(1)
    tests = [v for k, v in list(globals().items()) if k.startswith("test_")]
    failures = 0
    try:
        for t in tests:
            try:
                t()
                print(f"PASS {t.__name__}")
            except AssertionError as e:
                failures += 1
                print(f"FAIL {t.__name__}: {e}")
            except Exception as e:
                failures += 1
                print(f"ERROR {t.__name__}: {type(e).__name__}: {e}")
    finally:
        if _driver:
            _driver.close()
    print(f"\n{len(tests) - failures}/{len(tests)} passed")
    sys.exit(1 if failures else 0)
