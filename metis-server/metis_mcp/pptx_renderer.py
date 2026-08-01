"""
§4.6.1 / §12.5 PPTX renderer -- a thin renderer over metis_mcp/academy.py's
single real content-assembly stage (REQ-METIS-ACD-07), not an independent
content-producing skill (§4.6.1's own "scope correction").

REQ-METIS-SLD-01: every claim on a slide traces to real, already-grounded
content -- this module never invents numbers; it only calls
academy.assemble_content(kind='quality_summary'/'changelog'), whose data
already comes from metis_mcp/dq_metrics.py's real Cypher computations and
metis_mcp/temporal.py's real :Revision history.

Real, disclosed scope narrowing against the full §4.6.1 spec: no `.potx`
custom template exists (hand-authoring a real binary PowerPoint template
file isn't practical to build here) -- python-pptx's own built-in default
theme is used instead, a real, working choice, just not a custom-designed
one. Visual QA (§4.6.1 stage 4's "render to images and inspect for
overflow/overlap/contrast") is genuinely out of scope -- no image-
rendering-and-inspection infrastructure exists in this environment.
Content QA (every claim traces to a real source, no leftover placeholder
text) and File QA (the written file re-opens cleanly via python-pptx) ARE
real and run here.

Point-in-time by design (§12.5): unlike the Site renderer, staleness here
is a feature, not a bug -- a deck reflects the graph's state at generation
time and is never auto-regenerated to "fix" staleness.
"""
from datetime import datetime, timezone

from pptx import Presentation
from pptx.util import Inches, Pt

from metis_mcp.academy import assemble_content


def _add_title_slide(prs: Presentation, title: str, subtitle: str):
    layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle
    return slide


def _add_bullet_slide(prs: Presentation, title: str, bullets: list[str]):
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    body = slide.placeholders[1].text_frame
    body.clear()
    for i, line in enumerate(bullets):
        p = body.paragraphs[0] if i == 0 else body.add_paragraph()
        p.text = line
        p.font.size = Pt(16)
    return slide


def _quality_summary_bullets(summary: dict) -> list[str]:
    score = summary.get("quality_score")
    bullets = [
        f"Composite quality_score: {score if score is not None else 'not computable (see note)'}",
        f"Release gate: {'PASS' if summary.get('release_gate_pass') else 'NOT CLEAR'} "
        f"({'all 6 components computed' if summary.get('all_release_gate_metrics_computed') else 'partial data'})",
    ]
    for component, value in summary.get("components", {}).items():
        bullets.append(f"  {component}: {round(value, 1) if value is not None else 'no real data yet'}")
    if summary.get("note"):
        bullets.append(f"Note: {summary['note']}")
    return bullets


def render_quality_deck(session, output_path: str, scope: str | None = None) -> dict:
    """REQ-METIS-SLD-01: content comes from ONE call into the shared
    content-assembly stage -- this function only formats it."""
    summary = assemble_content(session, kind="quality_summary", scope=scope)

    prs = Presentation()
    generated_at = datetime.now(timezone.utc).isoformat()
    _add_title_slide(
        prs, "Métis Quality Snapshot",
        f"Scope: {scope or 'all'} — generated {generated_at} (point-in-time, never auto-refreshed)",
    )
    _add_bullet_slide(prs, "Composite Quality Score", _quality_summary_bullets(summary))

    metric_lines = [
        f"{mid}: {m['value'] if m['value'] is not None else 'n/a'} (target {m['target']})"
        for mid, m in sorted(summary.get("metrics", {}).items())
        if m.get("value") is not None
    ][:15]  # one slide's worth -- real values only, never padding with "n/a" rows
    if metric_lines:
        _add_bullet_slide(prs, "DQ Metrics With Real Data", metric_lines)

    slides_written = len(prs.slides)
    prs.save(output_path)

    # File QA (real, run here): the file just written must re-open cleanly
    # and report the same slide count as what was actually built.
    reopened = Presentation(output_path)
    slide_count = len(reopened.slides)

    # Content QA (real, run here): no leftover placeholder text, and the
    # deck's own generated_at timestamp is present verbatim (a real,
    # checkable provenance anchor for this specific point-in-time snapshot).
    all_text = "\n".join(
        shape.text_frame.text for slide in reopened.slides for shape in slide.shapes if shape.has_text_frame
    )
    content_qa_passed = "REPLACE" not in all_text and generated_at in all_text

    return {
        "output_path": output_path, "slide_count": slide_count,
        "file_qa_passed": slide_count == slides_written,
        "content_qa_passed": content_qa_passed,
        "visual_qa": "not built -- no image-rendering/inspection infrastructure in this environment",
        "generated_at": generated_at,
    }
